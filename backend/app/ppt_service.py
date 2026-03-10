from pathlib import Path
from uuid import uuid4

from pptx import Presentation


def build_pptx_file(title: str, slides: list[dict], output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    file_path = output_dir / f"lesson-{uuid4().hex}.pptx"

    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "AI 自动生成教学课件"

    body_layout = prs.slide_layouts[1]
    for slide in slides:
        page = prs.slides.add_slide(body_layout)
        page.shapes.title.text = slide.get("title", "未命名页面")
        text_frame = page.shapes.placeholders[1].text_frame
        text_frame.clear()

        points = slide.get("bullet_points", [])
        if not points:
            points = ["（本页暂无内容）"]

        for idx, point in enumerate(points):
            if idx == 0:
                text_frame.text = str(point)
            else:
                paragraph = text_frame.add_paragraph()
                paragraph.text = str(point)
                paragraph.level = 0

    prs.save(str(file_path))
    return file_path
